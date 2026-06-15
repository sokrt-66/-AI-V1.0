"""
多格式文档解析器：PDF / DOCX / TXT / MD / XLSX / CSV / PPTX
"""

import os
import re
import csv
import io
from typing import List, Dict

from ..core.logger import get_logger

logger = get_logger("rag.parser")


def parse_pdf(file_path: str) -> str:
    """解析 PDF 文档"""
    from PyPDF2 import PdfReader
    parts: List[str] = []
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text() or ""
            parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"PDF 解析异常: {e}")
        return ""


def parse_docx(file_path: str) -> str:
    """解析 Word 文档"""
    from docx import Document
    parts: List[str] = []
    try:
        doc = Document(file_path)
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"DOCX 解析异常: {e}")
        return ""


def parse_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.warning(f"TXT 解析异常: {e}")
        return ""


def parse_xlsx(file_path: str) -> str:
    """解析 Excel - 以 CSV 风格输出"""
    import openpyxl
    parts: List[str] = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"===== Sheet: {sheet_name} =====")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"XLSX 解析异常: {e}")
        return ""


def parse_csv(file_path: str) -> str:
    parts: List[str] = []
    try:
        with open(file_path, "r", encoding="utf-8-sig", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    parts.append(" | ".join(row))
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"CSV 解析异常: {e}")
        return ""


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    parts: List[str] = []
    try:
        prs = Presentation(file_path)
        for idx, slide in enumerate(prs.slides, 1):
            parts.append(f"===== 第 {idx} 页 =====")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"PPTX 解析异常: {e}")
        return ""


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".txt": parse_txt,
    ".md": parse_txt,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
    ".csv": parse_csv,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
}


def parse_document(file_path: str) -> Dict:
    """统一入口：解析文档并返回内容"""
    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSERS.get(ext)
    if not parser:
        # 尝试以文本方式读取
        logger.warning(f"未识别的文件类型 {ext}，尝试纯文本读取")
        parser = parse_txt

    raw = parser(file_path)
    from ..core.utils import clean_text
    content = clean_text(raw)

    return {
        "file_name": os.path.basename(file_path),
        "file_type": ext.lstrip("."),
        "file_size": os.path.getsize(file_path),
        "content": content,
        "content_length": len(content),
        "parser": parser.__name__,
    }
